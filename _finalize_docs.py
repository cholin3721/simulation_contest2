"""Update PPT, report, manual with re-simulated KPI (18,000s)."""
from pathlib import Path
import zipfile
import re
from pptx import Presentation

from _ppt_preserve_font import replace_in_shape_preserve

ROOT = Path(__file__).parent
PPT = ROOT / "K-Logistics_풀필먼트_센터_피킹_프로세스_효율화___인하이트_본선_발표_20260706095810.pptx"
MANUAL = ROOT / "제23회한국대학생컴퓨터시뮬레이션경진대회_인하이트사용자매뉴얼.docx"
REPORT = ROOT / "제23회한국대학생컴퓨터시뮬레이션경진대회_인하이트보고서.docx"

# AS-IS / TO-BE @ 18,000s (user screenshots)
ASIS = dict(pallets=8, boxes=256, rate=1.6, branches=1, agv=23, pct=2.9)
TOBE = dict(pallets=21, boxes=672, rate=4.2, branches=3, agv=50, pct=7.6)
INC_PCT = round((TOBE["boxes"] - ASIS["boxes"]) / ASIS["boxes"] * 100)

COMMON_REPLACEMENTS = [
    ("simulation_orders.xlsx", "database/db.script (SHEET1, 8,858박스)"),
    ("simulation_orders.xlsx  ← 모델(.alpx)과 동일 경로에 위치해야 함",
     "database/db.script (모델 내장 HSQLDB · SHEET1 8,858건)"),
    ("−transporter.distanceTo(agent)", "−unit.distanceTo(agent)"),
    ("transporter.distanceTo(agent)", "unit.distanceTo(agent)"),
    ("TransportTo 블록", "MoveByTransporter 블록"),
    ("TransportTo ·", "MoveByTransporter ·"),
    ("총 3개의 이동 블록", "총 5개의 MoveByTransporter 블록"),
    ("총 3개 이동 블록", "MoveByTransporter 5개 블록"),
    ("6,644 박스", "8,120 박스 (DB 실측 91.7%)"),
    ("6,644개", "8,120개 (91.7%)"),
    ("S/A 박스 6,644개(75%)", "S/A 박스 8,120개 (91.7%, DB)"),
    ("약 75%인 6,644 박스", "DB 실측 91.7%(8,120박스) · 이론 75%(6,644박스)"),
    ("S/A 박스 75%", "S/A 박스 91.7% (DB) · 이론 75%"),
    ("S/A 박스 6,644개", "S/A 박스 8,120개"),
]

KPI_REPLACEMENTS = [
    ("256 → 640 박스", f"{ASIS['boxes']} → {TOBE['boxes']} 박스"),
    ("256→640 박스", f"{ASIS['boxes']}→{TOBE['boxes']} 박스"),
    ("256 박스 → 640 박스", f"{ASIS['boxes']} 박스 → {TOBE['boxes']} 박스"),
    ("(+384박스)", f"(+{TOBE['boxes'] - ASIS['boxes']}박스)"),
    ("+384", f"+{TOBE['boxes'] - ASIS['boxes']}"),
    ("8→20 팔레트", f"{ASIS['pallets']}→{TOBE['pallets']} 팔레트"),
    ("8 → 20 팔레트", f"{ASIS['pallets']} → {TOBE['pallets']} 팔레트"),
    ("8 → 20 · 박스 처리 256 → 640", f"{ASIS['pallets']} → {TOBE['pallets']} · 박스 {ASIS['boxes']} → {TOBE['boxes']}"),
    ("완료 팔레트 8 → 20", f"완료 팔레트 {ASIS['pallets']} → {TOBE['pallets']}"),
    ("1.6→4.0 /hr", f"{ASIS['rate']}→{TOBE['rate']} /hr"),
    ("1.6 → 4.0 박스/hr", f"{ASIS['rate']} → {TOBE['rate']} 팔레트/hr"),
    ("1.6 → 4.0 /hr", f"{ASIS['rate']} → {TOBE['rate']} /hr"),
    ("1.6→4.0", f"{ASIS['rate']}→{TOBE['rate']}"),
    ("×2.5", "×2.6"),
    ("+150%", f"+{INC_PCT}%"),
    ("+150 %", f"+{INC_PCT} %"),
    ("1/60→3/60", f"{ASIS['branches']}/60→{TOBE['branches']}/60"),
    ("1 → 3", f"{ASIS['branches']} → {TOBE['branches']}"),
    ("1/60 → 3/60", f"{ASIS['branches']}/60 → {TOBE['branches']}/60"),
    ("23% → 44%", f"{ASIS['agv']}% → {TOBE['agv']}%"),
    ("약 23%", f"약 {ASIS['agv']}%"),
    ("약 44%", f"약 {TOBE['agv']}%"),
    ("23%→44%", f"{ASIS['agv']}%→{TOBE['agv']}%"),
    ("23%\n→\n44%", f"{ASIS['agv']}%\n→\n{TOBE['agv']}%"),
    ("×1.9", "×2.2"),
    ("2.9%", f"{ASIS['pct']}%"),
    ("7.2%", f"{TOBE['pct']}%"),
    ("25대 ★\n20 팔레트", f"25대 ★\n{TOBE['pallets']} 팔레트"),
    ("20 팔레트 (640 박스)", f"{TOBE['pallets']} 팔레트 ({TOBE['boxes']} 박스)"),
    ("20 팔레트 (640박스)", f"{TOBE['pallets']} 팔레트 ({TOBE['boxes']}박스)"),
    ("20 팔레트 · 4.0/hr", f"{TOBE['pallets']} 팔레트 · {TOBE['rate']}/hr"),
    ("약 44%", f"약 {TOBE['agv']}%"),
    ("8 팔레트 (256 박스)", f"{ASIS['pallets']} 팔레트 ({ASIS['boxes']} 박스)"),
    ("8 팔레트\n(256 박스)", f"{ASIS['pallets']} 팔레트\n({ASIS['boxes']} 박스)"),
    ("20 팔레트\n(640 박스)", f"{TOBE['pallets']} 팔레트\n({TOBE['boxes']} 박스)"),
    ("(+150%)", f"(+{INC_PCT}%)"),
    ("20 팔레트 (640 박스)", f"{TOBE['pallets']} 팔레트 ({TOBE['boxes']} 박스)"),
    ("8 팔레트\n256 박스", f"{ASIS['pallets']} 팔레트\n{ASIS['boxes']} 박스"),
    ("TO-BE\n25대 · 7.2%\n640 박스", f"TO-BE\n25대 · {TOBE['pct']}%\n{TOBE['boxes']} 박스"),
    ("AGV 가동률 23% → 44% (×1.9)", f"AGV 가동률 {ASIS['agv']}% → {TOBE['agv']}% (×{round(TOBE['agv']/ASIS['agv'],1)})"),
    ("완료 팔레트 8 → 20 (+150%)", f"완료 팔레트 {ASIS['pallets']} → {TOBE['pallets']} (+{INC_PCT}%)"),
    ("완료 지점 1/60 → 3/60 (×3배)", f"완료 지점 {ASIS['branches']}/60 → {TOBE['branches']}/60 (×3배)"),
    ("20 팔레트 (640 박스)", f"{TOBE['pallets']} 팔레트 ({TOBE['boxes']} 박스)"),
    ("8 팔레트 (256 박스)", f"{ASIS['pallets']} 팔레트 ({ASIS['boxes']} 박스)"),
    ("4.0/hr", f"{TOBE['rate']}/hr"),
    ("4.0 박스/hr", f"{TOBE['rate']} 팔레트/hr"),
    ("처리량 +150%", f"처리량 +{INC_PCT}%"),
    ("코드 변경 3곳으로 처리량 +150%", f"코드 변경 3곳으로 처리량 +{INC_PCT}%"),
]

# PPT-only: applied to whole-shape text (safe against run-splitting in docx).
PPT_EXTRA = [
    ("640 박스 (+416)", f"{TOBE['boxes']} 박스 (+{TOBE['boxes']-ASIS['boxes']})"),
    ("640 박스", f"{TOBE['boxes']} 박스"),
    ("20 팔레트 · 정상 · 기준", f"{TOBE['pallets']} 팔레트 · 정상 · 기준"),
    ("40대 (23%) → 25대 (44%)", f"40대 ({ASIS['agv']}%) → 25대 ({TOBE['agv']}%)"),
    ("제함기1 100% → 94%", "제함기1 100% → 약 30%"),
    ("제함기2 0% → 16%", "제함기2 0% → 약 70%"),
    ("완료 팔레트 8 → 20", f"완료 팔레트 {ASIS['pallets']} → {TOBE['pallets']}"),
    ("94%", "약 30%"),
    ("16%", "약 70%"),
    ("44%", f"{TOBE['agv']}%"),
]

REPORT_EXTRA = [
    ("20 팔레트 (640 박스)", f"{TOBE['pallets']} 팔레트 ({TOBE['boxes']} 박스)"),
    ("4.0/hr", f"{TOBE['rate']}/hr"),
    ("약 44%", f"약 {TOBE['agv']}%"),
    ("정적 슬롯 배정", "기본 슬롯 로직(selectOutput6) · TO-BE 동적 배정 강화"),
    ("파렛타이저 슬롯 배정 시", "selectOutput6/7 슬롯 배정 시"),
    ("TransportTo 블록의 TransporterRating", "MoveByTransporter 블록의 TransporterRating"),
    ("총 3개의 이동 블록에 적용", "총 5개의 MoveByTransporter 블록에 적용"),
    ("약 94%", "약 30%"),
    ("약16%", "약 70%"),
    ("25 (TO-BE)\n20 팔레트\n4.0/hr\n약 44%", f"25 (TO-BE)\n{TOBE['pallets']} 팔레트\n{TOBE['rate']}/hr\n약 {TOBE['agv']}%"),
    ("640 박스", f"{TOBE['boxes']} 박스"),
    ("7.2%", f"{TOBE['pct']}%"),
    ("simulation_orders.xlsx", "database/db.script (SHEET1)"),
]


def apply_replacements(text: str, pairs: list) -> str:
    for old, new in pairs:
        text = text.replace(old, new)
    return text


def replace_in_ppt(path: Path):
    prs = Presentation(path)
    pairs = COMMON_REPLACEMENTS + KPI_REPLACEMENTS + PPT_EXTRA
    n = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if replace_in_shape_preserve(shape, pairs):
                n += 1
    prs.save(path)
    return n


def replace_in_docx(path: Path, extra=None):
    pairs = COMMON_REPLACEMENTS + KPI_REPLACEMENTS + (extra or [])
    with zipfile.ZipFile(path, "r") as zin:
        buf = {}
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "word/document.xml":
                xml = data.decode("utf-8")
                for old, new in pairs:
                    xml = xml.replace(old, new)
                data = xml.encode("utf-8")
            buf[item.filename] = (item, data)
    with zipfile.ZipFile(path, "w") as zout:
        for fname, (item, data) in buf.items():
            zout.writestr(item, data)
    return len(pairs)


def main():
    import shutil
    ppt_paths = [PPT, PPT.with_suffix(".updated.pptx")]
    for ppt_target in ppt_paths:
        if not ppt_target.exists():
            continue
        try:
            n = replace_in_ppt(ppt_target)
            print(f"PPT {ppt_target.name}: {n} shapes")
        except PermissionError:
            tmp = PPT.with_suffix(".updated.pptx")
            if ppt_target != tmp:
                shutil.copy2(ppt_target, tmp)
            n = replace_in_ppt(tmp)
            print(f"PPT locked — wrote {tmp.name} ({n} shapes)")
    replace_in_docx(MANUAL)
    replace_in_docx(REPORT, REPORT_EXTRA)
    print("Manual + Report updated")
    print(f"KPI: AS-IS {ASIS} | TO-BE {TOBE} | +{INC_PCT}%")


if __name__ == "__main__":
    main()
