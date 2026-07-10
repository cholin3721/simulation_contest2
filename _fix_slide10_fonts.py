"""Re-apply slide 10 sensitivity data without breaking Pretendard formatting."""
from pathlib import Path

from pptx import Presentation

from _ppt_preserve_font import apply_table_map_preserve, replace_in_shape_preserve

ROOT = Path(__file__).parent
SRC = ROOT / "피킹_프로세스_효율화_및_생산성_극대화_20260709235709.pptx"
OUT = ROOT / "피킹_프로세스_효율화_및_생산성_극대화_20260709235709.updated.pptx"

TABLE_MAP = {
    (1, 1): "14 팔레트",
    (1, 2): "2.8 팔레트/hr",
    (1, 3): "약 80%",
    (1, 4): "후반 교착 · 운영 불가",
    (2, 1): "8 팔레트",
    (2, 2): "1.6 팔레트/hr",
    (2, 3): "약 35%",
    (2, 4): "1/3 지점 교착 · 운영 불가",
    (4, 1): "2 팔레트",
    (4, 2): "0.4 팔레트/hr",
    (4, 3): "약 30%",
    (4, 4): "초반 교착 · 운영 불가",
}

TEXT_REPLACEMENTS = [
    ("AGV 15 / 20 / 25 / 35 / 40 대", "AGV 15 / 20 / 25 / 35 대"),
    ("29 팔레트", "14 팔레트 · 후반 교착"),
    ("교착 · 측정 불가", "8 팔레트 · 1/3 교착"),
    ("4 팔레트 · 교착 발생", "2 팔레트 · 초반 교착"),
    (
        "15대는 정상이지만 처리량 여유가 부족하고 25대가 팔레트/hr · 가동률의 실용 최적점이다.",
        "15·20·35대는 교착으로 무진행 구간이 생기고, 25대만 5시간 끝까지 정상 운영되며 처리량·가동률의 실용 최적점이다.",
    ),
]


def main() -> None:
    prs = Presentation(SRC)
    slide = prs.slides[9]
    for shape in slide.shapes:
        if shape.has_table:
            apply_table_map_preserve(shape.table, TABLE_MAP)
        replace_in_shape_preserve(shape, TEXT_REPLACEMENTS)
    prs.save(OUT)
    print(f"saved {OUT.name} (fonts preserved)")


if __name__ == "__main__":
    main()
