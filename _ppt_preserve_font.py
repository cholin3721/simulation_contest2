"""PPT text edits that keep existing run formatting (font, size, bold, color)."""
from __future__ import annotations

from pptx.oxml.ns import qn


def paragraph_full_text(paragraph) -> str:
    return "".join(run.text for run in paragraph.runs)


def set_paragraph_text_preserve(paragraph, new_text: str) -> None:
    """Replace paragraph text without resetting the first run's rPr."""
    p = paragraph._p
    runs = p.findall(qn("a:r"))
    if not runs:
        paragraph.text = new_text
        return
    t = runs[0].find(qn("a:t"))
    if t is None:
        from pptx.oxml import parse_xml

        runs[0].append(parse_xml(f'<a:t xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">{new_text}</a:t>'))
    else:
        t.text = new_text
    for r in runs[1:]:
        p.remove(r)


def set_text_frame_text_preserve(text_frame, new_text: str) -> None:
    if not text_frame.paragraphs:
        text_frame.text = new_text
        return
    set_paragraph_text_preserve(text_frame.paragraphs[0], new_text)
    for p in text_frame.paragraphs[1:]:
        set_paragraph_text_preserve(p, "")


def set_cell_text_preserve(cell, new_text: str) -> None:
    set_text_frame_text_preserve(cell.text_frame, new_text)


def replace_in_paragraph_preserve(paragraph, replacements: list[tuple[str, str]]) -> bool:
    full = paragraph_full_text(paragraph)
    updated = full
    for old, new in replacements:
        updated = updated.replace(old, new)
    if updated == full:
        return False
    set_paragraph_text_preserve(paragraph, updated)
    return True


def replace_in_text_frame_preserve(text_frame, replacements: list[tuple[str, str]]) -> bool:
    changed = False
    for p in text_frame.paragraphs:
        if replace_in_paragraph_preserve(p, replacements):
            changed = True
    return changed


def replace_in_shape_preserve(shape, replacements: list[tuple[str, str]]) -> bool:
    changed = False
    if getattr(shape, "has_table", False) and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if replace_in_text_frame_preserve(cell.text_frame, replacements):
                    changed = True
    elif getattr(shape, "has_text_frame", False) and shape.has_text_frame:
        if replace_in_text_frame_preserve(shape.text_frame, replacements):
            changed = True
    return changed


def apply_table_map_preserve(table, cell_map: dict[tuple[int, int], str]) -> None:
    for (r, c), value in cell_map.items():
        set_cell_text_preserve(table.cell(r, c), value)
